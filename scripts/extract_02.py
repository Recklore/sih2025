import os
import shutil
import subprocess
from typing import Optional

import fitz  # PyMuPDF
from bs4 import BeautifulSoup
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def is_digital(pdf_path: str) -> bool:
    try:
        with fitz.open(pdf_path) as doc:
            text_pages = sum(1 for page in doc if page.get_text().strip())
            if doc.page_count / 2 > text_pages:
                return False
            return True
    except Exception:
        return False


def ghostscript_repair(pdf_path: str, repaired_dir: str) -> str:
    os.makedirs(repaired_dir, exist_ok=True)
    out_path = os.path.join(repaired_dir, os.path.basename(pdf_path))

    # Try common Ghostscript executables on Windows, then generic 'gs'
    gs_candidates = ["gswin64c", "gswin32c", "gs"]
    gs = None
    for cand in gs_candidates:
        try:
            subprocess.run([cand, "-v"], check=True, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
            gs = cand
            break
        except Exception:
            continue

    if gs is None:
        raise RuntimeError("Ghostscript not found. Install it and ensure 'gswin64c' is on PATH.")

    subprocess.run(
        [gs, "-o", out_path, "-sDEVICE=pdfwrite", "-dPDFSETTINGS=/prepress", pdf_path],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )
    return out_path


def ocr_extract(pdf_path: str, ocr_dir: str, reader, dpi: int = 200) -> str:
    os.makedirs(ocr_dir, exist_ok=True)
    base = os.path.splitext(os.path.basename(pdf_path))[0]
    txt_path = os.path.join(ocr_dir, f"{base}_ocr.txt")

    with fitz.open(pdf_path) as doc, open(txt_path, "w", encoding="utf-8") as txt_file:
        for page_num in range(doc.page_count):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=dpi)
            # Following the existing pattern in your repo: feed bytes to EasyOCR
            ocr_result = reader.readtext(pix.tobytes(), detail=0, paragraph=True)
            txt_file.write(f"Page {page_num + 1}:\n")
            txt_file.write(" ".join(ocr_result))
            txt_file.write("\n" + "-" * 80 + "\n")

    return txt_path


def extract_txt(file_path: str) -> str:
    """Extract text from plain text file."""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        # Fallback to latin-1 encoding
        with open(file_path, "r", encoding="latin-1") as f:
            return f.read()


def extract_docx(file_path: str) -> str:
    """Extract text from Word document."""
    doc = Document(file_path)
    text_parts = []

    # Extract paragraphs
    for para in doc.paragraphs:
        if para.text.strip():
            text_parts.append(para.text)

    # Extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells)
            if row_text.strip():
                text_parts.append(row_text)

    return "\n".join(text_parts)


def extract_xlsx(file_path: str) -> str:
    """Extract text from Excel spreadsheet."""
    wb = load_workbook(file_path, read_only=True, data_only=True)
    text_parts = []

    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        text_parts.append(f"\n=== Sheet: {sheet_name} ===\n")

        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                text_parts.append(row_text)

    wb.close()
    return "\n".join(text_parts)


def extract_pptx(file_path: str) -> str:
    """Extract text from PowerPoint presentation."""
    prs = Presentation(file_path)
    text_parts = []

    for slide_num, slide in enumerate(prs.slides, 1):
        text_parts.append(f"\n=== Slide {slide_num} ===\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text)

    return "\n".join(text_parts)


def extract_html(file_path: str) -> str:
    """Extract text from HTML file."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f, "lxml")

    # Remove script and style elements
    for script in soup(["script", "style", "nav", "footer", "header"]):
        script.decompose()

    # Get text
    text = soup.get_text(separator="\n", strip=True)

    # Clean up whitespace
    lines = (line.strip() for line in text.splitlines())
    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
    text = "\n".join(chunk for chunk in chunks if chunk)

    return text


def extract_pdf(file_path: str, output_dir: str, reader=None) -> str:
    """
    Extract text from PDF (digital or scanned).
    Returns path to saved text file.
    """
    os.makedirs(output_dir, exist_ok=True)
    base = os.path.splitext(os.path.basename(file_path))[0]
    txt_path = os.path.join(output_dir, f"{base}.txt")

    try:
        # Check if digital PDF
        if is_digital(file_path):
            # Extract text directly
            with fitz.open(file_path) as doc:
                text_parts = []
                for page_num, page in enumerate(doc, 1):
                    text = page.get_text()
                    if text.strip():
                        text_parts.append(f"Page {page_num}:\n{text}")
                        text_parts.append("-" * 80)

                full_text = "\n".join(text_parts)

            # Clean the extracted text
            full_text = clean_text(full_text)

            with open(txt_path, "w", encoding="utf-8") as f:
                f.write(full_text)

            return txt_path
        else:
            # Scanned PDF - use OCR
            if reader is None:
                import easyocr

                reader = easyocr.Reader(["hi", "en"])

            with fitz.open(file_path) as doc:
                text_parts = []
                for page_num in range(doc.page_count):
                    page = doc.load_page(page_num)
                    pix = page.get_pixmap(dpi=200)
                    ocr_result = reader.readtext(pix.tobytes(), detail=0, paragraph=True)
                    text_parts.append(f"Page {page_num + 1}:\n")
                    text_parts.append(" ".join(ocr_result))
                    text_parts.append("-" * 80)

                full_text = "\n".join(text_parts)

            # Clean the extracted text
            full_text = clean_text(full_text)

            with open(txt_path, "w", encoding="utf-8") as f:
                f.write(full_text)

            return txt_path

    except Exception as e:
        raise Exception(f"Failed to extract PDF: {e}")


def clean_text(text: str) -> str:
    """
    Clean extracted text by removing page separators and normalizing whitespace.
    Replaces page markers (--- or 'Page N:') with # symbols.
    Removes excessive empty lines while preserving structure.
    """
    lines = text.split("\n")
    cleaned_lines = []
    is_header_block = False

    for line in lines:
        stripped_line = line.strip()

        # Detect page separators or page headers
        if "---" in stripped_line or stripped_line.startswith("Page "):
            if not is_header_block:
                cleaned_lines.append("#")
                is_header_block = True

        # Keep non-empty content lines
        elif stripped_line:
            cleaned_lines.append(line)
            is_header_block = False

    return "\n".join(cleaned_lines)


def process_file(file_path: str, output_dir: str, reader=None) -> Optional[str]:
    """
    Process a single file and extract text.
    Returns path to output text file or None on failure.
    """
    file_ext = os.path.splitext(file_path)[1].lower()
    base_name = os.path.splitext(os.path.basename(file_path))[0]
    output_path = os.path.join(output_dir, f"{base_name}.txt")

    try:
        os.makedirs(output_dir, exist_ok=True)

        # Route to appropriate extractor
        if file_ext == ".pdf":
            return extract_pdf(file_path, output_dir, reader)

        elif file_ext == ".txt":
            text = extract_txt(file_path)

        elif file_ext in [".doc", ".docx"]:
            text = extract_docx(file_path)

        elif file_ext in [".xls", ".xlsx"]:
            text = extract_xlsx(file_path)

        elif file_ext in [".ppt", ".pptx"]:
            text = extract_pptx(file_path)

        elif file_ext in [".html", ".htm"]:
            text = extract_html(file_path)

        else:
            print(f"Unsupported file type: {file_ext}")
            return None

        # Clean the extracted text
        text = clean_text(text)

        # Save extracted text
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(text)

        return output_path

    except Exception as e:
        print(f"Error processing {os.path.basename(file_path)}: {e}")
        return None


def process_pdf(pdf_path: str, repaired_dir: str, ocr_dir: str, error_dir: Optional[str] = None, reader=None) -> None:
    try:
        if is_digital(pdf_path):
            ghostscript_repair(pdf_path, repaired_dir)
            print(f"Digital (repaired): {os.path.basename(pdf_path)}")
        else:
            # Lazy import to avoid requiring EasyOCR if only repairing
            global easyocr
            if reader is None:
                import easyocr  # type: ignore

                reader = easyocr.Reader(["hi", "en"])  # initialize once if needed
            ocr_extract(pdf_path, ocr_dir, reader)
            print(f"Scanned (OCRed): {os.path.basename(pdf_path)}")
    except Exception as e:
        print(f"Error processing {os.path.basename(pdf_path)}: {e}")
        if error_dir:
            os.makedirs(error_dir, exist_ok=True)
            try:
                shutil.copy(pdf_path, os.path.join(error_dir, os.path.basename(pdf_path)))
            except Exception:
                pass


def main(
    input_base_dir: str = "./data",
    output_base_dir: str = "./processed_data",
    limit_per_type: Optional[int] = 50,
):
    """
    Process files from data/{pdf,docs,html}/ and extract text to processed_data/{pdf,docs,html}/.
    """
    categories = ["pdf", "docs", "html"]

    # Initialize EasyOCR reader once for scanned PDFs
    reader = None

    total_processed = 0
    total_failed = 0

    for category in categories:
        input_dir = os.path.join(input_base_dir, category)
        output_dir = os.path.join(output_base_dir, category)

        if not os.path.isdir(input_dir):
            print(f"Input directory '{input_dir}' does not exist, skipping...")
            continue

        os.makedirs(output_dir, exist_ok=True)

        print(f"\n{'='*60}")
        print(f"Processing {category.upper()} files from: {input_dir}")
        print(f"{'='*60}")

        files = sorted(os.listdir(input_dir))
        count = 0

        for filename in files:
            file_path = os.path.join(input_dir, filename)

            # Skip directories
            if not os.path.isfile(file_path):
                continue

            # Check limit
            if limit_per_type is not None and count >= limit_per_type:
                print(f"Reached limit of {limit_per_type} files for {category}")
                break

            # Initialize EasyOCR only when needed (lazy loading)
            if category == "pdf" and reader is None:
                try:
                    import easyocr

                    print("Initializing EasyOCR (this may take a moment)...")
                    reader = easyocr.Reader(["hi", "en"], gpu=True)
                except Exception as e:
                    print(f"Warning: EasyOCR initialization failed: {e}")
                    reader = None

            # Process file
            result = process_file(file_path, output_dir, reader)

            if result:
                print(f"✅ {filename} → {os.path.basename(result)}")
                total_processed += 1
            else:
                print(f"❌ {filename} - Failed")
                total_failed += 1

            count += 1

    print(f"\n{'='*60}")
    print(f"Processing completed!")
    print(f"Total files processed: {total_processed}")
    print(f"Total files failed: {total_failed}")
    print(f"Output directory: {output_base_dir}")
    print(f"{'='*60}")


if __name__ == "__main__":
    main()
